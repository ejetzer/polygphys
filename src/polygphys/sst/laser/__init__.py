#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Créer de nouveaux certificats de sureté laser.

Created on Thu Mar  3 11:22:01 2022

@author: emilejetzer
"""

import subprocess
import getpass
import time

import tkinter as tk
import configparser as cp

from pathlib import Path
from datetime import datetime as dt
from subprocess import run

import pptx
import pandas as pd

from ...outils.reseau.msforms import MSFormConfig, MSForm
from ...outils.reseau import FichierLointain, DisqueRéseau


class SSTLaserCertificatsConfig(MSFormConfig):

    def default(self):
        return (Path(__file__).parent / 'nouveau_certificat.cfg').open().read()


class SSTLaserCertificatsForm(MSForm):

    def nettoyer(self, cadre):
        cadre = self.convertir_champs(cadre)
        cadre = cadre.astype({'matricule': int}, errors='ignore')

        courriels_manquants = cadre['courriel'] == 'anonymous'
        cadre.loc[courriels_manquants,
                  'courriel'] = cadre.loc[courriels_manquants, 'courriel2']
        cadre.courriel = cadre.courriel.fillna(
            cadre.courriel2).fillna('@polymtl.ca')
        cadre.nom = cadre.nom.fillna(cadre.nom2).fillna('anonyme')
        cadre.date = cadre.date.dt.date
        cadre.matricule = cadre.matricule.fillna(0)

        return cadre.loc[:, ['date', 'matricule', 'courriel', 'nom']]

    def action(self, cadre):
        for i, entrée in cadre.iterrows():
            chemin_cert = Path(__file__).parent / \
                self.config.get('certificats', 'chemin')
            cert = pptx.Presentation(chemin_cert)

            for forme in cert.slides[0].shapes:
                if forme.has_text_frame:
                    for par in forme.text_frame.paragraphs:
                        for ligne in par.runs:
                            if ligne.text == 'nom':
                                ligne.text = str(entrée.nom)
                            elif ligne.text == 'matricule':
                                ligne.text = str(entrée.matricule)
                            elif ligne.text.startswith('Date'):
                                date = dt.today()
                                ligne.text = f'Date: {date.year}-{date.month:02}'

            for disque in self.config.getlist('certificats', 'disques'):
                url = self.config.get(disque, 'url')
                chemin = self.config.getpath(disque, 'mount_point')
                drive = self.config.get(disque, 'drive')

                print(f'Connection à {disque}')
                nom = input('nom: ')
                mdp = getpass.getpass('mdp: ')
                with DisqueRéseau(url, chemin, drive, nom, mdp) as d:
                    sous_dossier = d / self.config.get(disque, 'chemin')
                    sous_dossier = d / self.config.get('certificats', 'ppt')
                    fichier = sous_dossier / f'{entrée.nom}.pptx'
                    cert.save(fichier)

                    fichier_pdf = fichier.parent.parent / 'pdf' / fichier.name
                    run(['unoconv',
                         '-f',
                         'pdf',
                         '-o',
                         str(fichier_pdf),
                         str(fichier)])


def main():
    chemin_config = next(Path(__file__).parent.glob('*.cfg'))
    config = SSTLaserCertificatsConfig(chemin_config)

    fichier = FichierLointain(config.get('formulaire', 'url'),
                              config.getpath('formulaire', 'chemin'))
    fichier.update()

    formulaire = SSTLaserCertificatsForm(config)

    exporteur = subprocess.Popen(['unoconv', '--listener'])

    try:
        while True:
            fichier.update()
            formulaire.mise_à_jour()
            time.sleep(60 * 60)  # On roule à chaque heure
    finally:
        config.getpath('formulaire', 'chemin').unlink()
        exporteur.terminate()


if __name__ == '__main__':
    main()
